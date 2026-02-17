"""Impress-category fixture file generators (pptx, odp)."""

from io import BytesIO
from pathlib import Path
from typing import Optional


# ---------------------------------------------------------------------------
# Shared test data — "분기 실적 발표" scenario (FR-018)
# ---------------------------------------------------------------------------

IMPRESS_COMPANY = "LibreFormer Inc."
IMPRESS_DATE = "2026 Q1"

IMPRESS_TITLE_SLIDE = {
    "title": IMPRESS_COMPANY,
    "subtitle": IMPRESS_DATE,
}

IMPRESS_CONTENT_SLIDE = {
    "title": "핵심 지표",
    "bullets": [
        "매출 성장률: 25%",
        "신규 고객 수: 150개사",
        "고객 만족도: 4.8/5.0",
        "서비스 가동률: 99.9%",
    ],
}

IMPRESS_IMAGE_SLIDE = {
    "title": "실적 차트",
}


# ---------------------------------------------------------------------------
# T022: create_pptx
# ---------------------------------------------------------------------------


def create_pptx(path: Path, image_bytes: Optional[bytes] = None) -> Path:
    """Create a .pptx file with title slide, content slide, and image slide.

    Uses the "분기 실적 발표" (Quarterly Performance Report) scenario data (FR-018).

    Args:
        path: Destination file path.
        image_bytes: Optional PNG image bytes for the image slide.

    Returns:
        The path to the created .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # --- Slide 1: Title slide ---
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = IMPRESS_TITLE_SLIDE["title"]
    slide.placeholders[1].text = IMPRESS_TITLE_SLIDE["subtitle"]

    # --- Slide 2: Content slide with bullet points ---
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = IMPRESS_CONTENT_SLIDE["title"]
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(IMPRESS_CONTENT_SLIDE["bullets"]):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet

    # --- Slide 3: Image slide ---
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    # Add title as a text box

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = IMPRESS_IMAGE_SLIDE["title"]
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    if image_bytes:
        img_stream = BytesIO(image_bytes)
        slide.shapes.add_picture(
            img_stream, Inches(2.5), Inches(1.5), Inches(5), Inches(5)
        )

    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# T023: create_odp
# ---------------------------------------------------------------------------


def create_odp(path: Path) -> Path:
    """Create a .odp file with title and content slides using odfpy.

    Uses the "분기 실적 발표" scenario data (FR-018).

    Args:
        path: Destination file path.

    Returns:
        The path to the created .odp file.
    """
    import pytest

    pytest.importorskip("odf")

    from odf.opendocument import OpenDocumentPresentation
    from odf.style import Style, MasterPage, PageLayout, PageLayoutProperties
    from odf.style import (
        GraphicProperties,
        DrawingPageProperties,
        TextProperties,
    )
    from odf.draw import Page, Frame, TextBox
    from odf.text import P

    doc = OpenDocumentPresentation()

    # Page layout
    pagelayout = PageLayout(name="MyLayout")
    doc.automaticstyles.addElement(pagelayout)
    pagelayout.addElement(
        PageLayoutProperties(
            margin="0cm",
            pagewidth="28cm",
            pageheight="21cm",
            printorientation="landscape",
        )
    )

    # Master page
    masterpage = MasterPage(name="MyMaster", pagelayoutname=pagelayout)
    doc.masterstyles.addElement(masterpage)

    # Drawing page style
    dpstyle = Style(name="dp1", family="drawing-page")
    dpstyle.addElement(DrawingPageProperties(transitiontype="none"))
    doc.automaticstyles.addElement(dpstyle)

    # Presentation style for text
    presstyle = Style(name="pr1", family="presentation")
    presstyle.addElement(GraphicProperties(fillcolor="#ffffff"))
    doc.automaticstyles.addElement(presstyle)

    # Title style
    titlestyle = Style(name="title", family="paragraph")
    titlestyle.addElement(TextProperties(fontsize="32pt", fontweight="bold"))
    doc.automaticstyles.addElement(titlestyle)

    # Subtitle style
    substyle = Style(name="subtitle", family="paragraph")
    substyle.addElement(TextProperties(fontsize="20pt"))
    doc.automaticstyles.addElement(substyle)

    # Bullet style
    bulletstyle = Style(name="bullet", family="paragraph")
    bulletstyle.addElement(TextProperties(fontsize="18pt"))
    doc.automaticstyles.addElement(bulletstyle)

    # --- Slide 1: Title ---
    page = Page(stylename=dpstyle, masterpagename=masterpage)
    doc.presentation.addElement(page)

    frame = Frame(stylename=presstyle, width="25cm", height="6cm", x="1.5cm", y="5cm")
    page.addElement(frame)
    textbox = TextBox()
    frame.addElement(textbox)
    title_p = P(stylename=titlestyle)
    title_p.addText(IMPRESS_TITLE_SLIDE["title"])
    textbox.addElement(title_p)
    sub_p = P(stylename=substyle)
    sub_p.addText(IMPRESS_TITLE_SLIDE["subtitle"])
    textbox.addElement(sub_p)

    # --- Slide 2: Content with bullets ---
    page2 = Page(stylename=dpstyle, masterpagename=masterpage)
    doc.presentation.addElement(page2)

    frame2 = Frame(stylename=presstyle, width="25cm", height="18cm", x="1.5cm", y="1cm")
    page2.addElement(frame2)
    textbox2 = TextBox()
    frame2.addElement(textbox2)
    content_title = P(stylename=titlestyle)
    content_title.addText(IMPRESS_CONTENT_SLIDE["title"])
    textbox2.addElement(content_title)
    for bullet in IMPRESS_CONTENT_SLIDE["bullets"]:
        bp = P(stylename=bulletstyle)
        bp.addText(f"• {bullet}")
        textbox2.addElement(bp)

    doc.save(str(path))
    return path


# ---------------------------------------------------------------------------
# T028: create_empty_pptx (edge case)
# ---------------------------------------------------------------------------


def create_empty_pptx(path: Path) -> Path:
    """Create an empty .pptx file with no slides.

    Args:
        path: Destination file path.

    Returns:
        The path to the created .pptx file.
    """
    from pptx import Presentation

    prs = Presentation()
    # No slides added — empty presentation
    prs.save(str(path))
    return path
